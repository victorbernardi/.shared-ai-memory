# coding: utf-8
"""Extrai imagens embutidas dos slides 6 e 7 do PPTX."""
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = Path(r"C:\Projetos\Inova\projects\lead-csc-pops\docs\business\Campanha-Leads-Preventivos-Inova.pptx")
OUT = Path(r"C:\Projetos\Inova\projects\lead-csc-pops\docs\business")

prs = Presentation(str(PPTX))
for si in [5, 6]:  # slides 6 e 7 (0-indexed)
    slide = prs.slides[si]
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            name = shape.name
            img = shape.image
            ext = img.ext
            out_file = OUT / f"_slide{si+1}_{name}.{ext}"
            out_file.write_bytes(img.blob)
            print(f"Extracted: {out_file.name} ({len(img.blob):,} bytes)")
