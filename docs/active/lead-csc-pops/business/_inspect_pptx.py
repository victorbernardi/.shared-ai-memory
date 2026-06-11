# coding: utf-8
"""
_inspect_pptx.py
Varre o PPTX e imprime a estrutura completa de cada slide:
shapes, tipo, posição, tamanho, textos, cores, fontes.
"""
import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = Path(r"C:\Projetos\Inova\projects\lead-csc-pops\docs\business\Campanha-Leads-Preventivos-Inova.pptx")

def emu_to_in(emu):
    return round(emu / 914400, 4)

def rgb_str(color):
    try:
        return f"#{color.rgb}"
    except:
        return "inherited/theme"

def print_tf(tf, indent=6):
    for i, para in enumerate(tf.paragraphs):
        text = para.text
        if not text.strip():
            continue
        sp = " " * indent
        align = str(para.alignment) if para.alignment else "None"
        sb = para.space_before.pt if para.space_before else 0
        print(f"{sp}[PARA {i}] align={align} space_before={sb:.0f}pt")
        print(f"{sp}  TEXT: {repr(text[:120])}")
        for run in para.runs:
            fn = run.font
            size = f"{fn.size.pt:.1f}pt" if fn.size else "None"
            bold = fn.bold
            color = rgb_str(fn.color) if fn.color and fn.color.type else "None"
            name = fn.name or "None"
            print(f"{sp}  RUN: font={name} size={size} bold={bold} color={color} text={repr(run.text[:80])}")

def inspect():
    prs = Presentation(str(PPTX))
    w = emu_to_in(prs.slide_width)
    h = emu_to_in(prs.slide_height)
    print(f"SLIDE SIZE: {w}\" x {h}\"")
    print(f"TOTAL SLIDES: {len(prs.slides)}\n")
    print("=" * 80)

    for si, slide in enumerate(prs.slides):
        print(f"\n{'=' * 80}")
        print(f"SLIDE {si + 1}")
        print(f"{'=' * 80}")

        for shi, shape in enumerate(slide.shapes):
            l = emu_to_in(shape.left) if shape.left is not None else "?"
            t = emu_to_in(shape.top) if shape.top is not None else "?"
            w = emu_to_in(shape.width) if shape.width is not None else "?"
            h = emu_to_in(shape.height) if shape.height is not None else "?"
            stype = str(shape.shape_type)

            print(f"\n  [{shi}] {shape.shape_type.name if hasattr(shape.shape_type, 'name') else stype}")
            print(f"       name={repr(shape.name)}")
            print(f"       pos=({l}\", {t}\") size=({w}\" x {h}\")")

            # Fill
            try:
                fill = shape.fill
                if fill.type is not None:
                    try:
                        fc = f"#{fill.fore_color.rgb}"
                    except:
                        fc = "N/A"
                    print(f"       fill_type={fill.type} fore_color={fc}")
            except:
                pass

            # Line
            try:
                line = shape.line
                if line.color and line.color.type:
                    lc = f"#{line.color.rgb}"
                    lw = line.width.pt if line.width else "?"
                    print(f"       line_color={lc} line_width={lw}pt")
            except:
                pass

            # Text
            if shape.has_text_frame:
                print(f"       HAS TEXT:")
                print_tf(shape.text_frame, indent=8)

            # Picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"       PICTURE: {getattr(shape, 'name', '?')}")

            # Table
            if shape.has_table:
                tbl = shape.table
                print(f"       TABLE: {tbl.rows.__len__()} rows x {len(tbl.columns)} cols")
                for ri, row in enumerate(tbl.rows):
                    for ci, cell in enumerate(row.cells):
                        ct = cell.text_frame.text.strip()
                        if ct:
                            try:
                                bg = f"#{cell.fill.fore_color.rgb}"
                            except:
                                bg = "N/A"
                            print(f"         [{ri},{ci}] bg={bg} text={repr(ct[:60])}")

if __name__ == "__main__":
    inspect()
