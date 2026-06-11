# coding: utf-8
"""
generate_slides_leads.py
Gera apresentação PPTX premium da Campanha de Leads Preventivos de Peças
(FPS e Material Rodante) para a Inova Máquinas.

Layout: dark brutalist reaproveitado de apresentacao-roberto-2505.
"""

import os
import subprocess
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
BASE_DIR = Path(__file__).parent
TEMPLATE_DIR = Path(r"C:\Projetos\Inova\projects\lead-csc-pops\Template")
LOGO_INOVA = TEMPLATE_DIR / "inova_logo_cdcdcd.png"
LOGO_JD = TEMPLATE_DIR / "JD_logo_no_text.png"
KPI_HTML = Path(r"C:\Projetos\Inova\projects\lead-csc-pops\data\output\daily_report_kpis.html")
KPI_SCREENSHOT = BASE_DIR / "_kpi_screenshot.png"
OUTPUT_PPTX = BASE_DIR / "Campanha-Leads-Preventivos-Inova.pptx"

# ---------------------------------------------------------------------------
# Theme Colors (idêntico ao projeto roberto-2505)
# ---------------------------------------------------------------------------
COLOR_DARK_BG   = RGBColor(17, 17, 17)
COLOR_YELLOW_JD = RGBColor(255, 194, 14)
COLOR_GRAY_DARK = RGBColor(31, 41, 55)
COLOR_GRAY_TEXT = RGBColor(156, 163, 175)
COLOR_WHITE     = RGBColor(255, 255, 255)
COLOR_GREEN     = RGBColor(34, 197, 94)
COLOR_RED       = RGBColor(239, 68, 68)
COLOR_AMBER     = RGBColor(245, 158, 11)


# ---------------------------------------------------------------------------
# Core layout helpers (replicados do projeto de referência)
# ---------------------------------------------------------------------------

def apply_slide_theme(prs: Presentation, slide, is_cover: bool = False) -> None:
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DARK_BG

    # Barra vertical amarela (esquerda)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.15), Inches(0), Inches(0.15), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_YELLOW_JD
    bar.line.fill.background()

    if not is_cover:
        # Banner cinza topo
        top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.4))
        top.fill.solid()
        top.fill.fore_color.rgb = COLOR_GRAY_DARK
        top.line.fill.background()

        # Linha rodapé
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), prs.slide_height - Inches(0.1), prs.slide_width, Inches(0.1))
        footer.fill.solid()
        footer.fill.fore_color.rgb = COLOR_GRAY_DARK
        footer.line.fill.background()

        # Logos no header (topo direito)
        try:
            if LOGO_INOVA.exists():
                slide.shapes.add_picture(str(LOGO_INOVA), Inches(10.8), Inches(0.06), Inches(1.0), Inches(0.28))
            if LOGO_JD.exists():
                slide.shapes.add_picture(str(LOGO_JD), Inches(12.0), Inches(0.06), Inches(0.58), Inches(0.28))
        except Exception as e:
            print(f"[WARNING] Logos: {e}")


def create_eyebrow(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = text.upper()
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_YELLOW_JD
    p.font.bold = True


def create_slide_title(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.5), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True


def add_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_GRAY_DARK
    card.line.color.rgb = COLOR_YELLOW_JD
    card.line.width = Pt(1.5)
    return card


def add_callout(slide, text: str, color: RGBColor = None) -> None:
    add_card(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8))
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(6.42), Inches(11.3), Inches(0.76))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(10.5)
    p.font.color.rgb = color or COLOR_YELLOW_JD
    p.font.bold = True


def add_text_block(slide, left, top, width, height, lines: list[tuple]) -> None:
    """
    lines: lista de (texto, tamanho_pt, RGBColor, bold, space_before_pt)
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    first = True
    for (text, size, color, bold, space) in lines:
        if first:
            p = tb.text_frame.paragraphs[0]
            first = False
        else:
            p = tb.text_frame.add_paragraph()
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        if space:
            p.space_before = Pt(space)


# ---------------------------------------------------------------------------
# Screenshot do HTML de KPIs via Edge headless
# ---------------------------------------------------------------------------

def capture_kpi_screenshot() -> bool:
    if KPI_SCREENSHOT.exists():
        print(f" -> Usando screenshot existente: {KPI_SCREENSHOT.name}")
        return True
    edge = r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"
    if not os.path.exists(edge):
        print("[WARNING] Edge não encontrado. Slides de KPI sem imagem.")
        return False
    result = subprocess.run([
        edge, "--headless", "--disable-gpu", "--no-sandbox",
        "--screenshot=" + str(KPI_SCREENSHOT),
        "--window-size=1400,900",
        f"file:///{KPI_HTML}",
    ], capture_output=True, timeout=30)
    if KPI_SCREENSHOT.exists():
        print(f" -> Screenshot gerado: {KPI_SCREENSHOT.name}")
        return True
    print(f"[WARNING] Screenshot falhou: {result.stderr[:200]}")
    return False


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------

def build_slide_capa(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide, is_cover=True)

    # Logos centralizados na capa
    try:
        if LOGO_INOVA.exists():
            slide.shapes.add_picture(str(LOGO_INOVA), Inches(4.5), Inches(1.3), Inches(1.8), Inches(0.5))
        if LOGO_JD.exists():
            slide.shapes.add_picture(str(LOGO_JD), Inches(6.7), Inches(1.2), Inches(1.0), Inches(0.6))
    except Exception as e:
        print(f"[WARNING] Logos capa: {e}")

    # Título principal
    tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0))
    p = tb.text_frame.paragraphs[0]
    p.text = "CAMPANHA DE LEADS PREVENTIVOS DE PEÇAS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(36)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    p2 = tb.text_frame.add_paragraph()
    p2.text = "FPS & Material Rodante — Inova Máquinas"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_YELLOW_JD
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(14)

    # Linha decorativa
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.6), Inches(6.3), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_YELLOW_JD
    line.line.fill.background()

    # Descrição
    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(0.6))
    p3 = tb2.text_frame.paragraphs[0]
    p3.text = "Motor automatizado de geração de leads com base em horímetros de máquinas John Deere"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_GRAY_TEXT
    p3.alignment = PP_ALIGN.CENTER


def build_slide_problema(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "diagnóstico comercial // peças")
    create_slide_title(slide, "O CUSTO DO PROCESSO REATIVO")

    # Card esquerdo — situação atual
    add_card(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.7))
    add_text_block(slide, Inches(1.0), Inches(1.65), Inches(5.2), Inches(4.4), [
        ("SITUAÇÃO ATUAL", 10, COLOR_GRAY_TEXT, True, 0),
        ("O cliente liga quando a peça quebra — ou não liga e vai ao concorrente.", 10.5, COLOR_WHITE, False, 8),
        ("• Nenhum consultor recebe alerta automático de desgaste", 10, COLOR_WHITE, False, 10),
        ("• O contato depende da memória do vendedor ou do cliente", 10, COLOR_WHITE, False, 6),
        ("• Peças compradas em emergência vão para o fornecedor mais rápido", 10, COLOR_WHITE, False, 6),
        ("• Não há como auditar se o contato foi feito e gerou proposta", 10, COLOR_WHITE, False, 6),
        ("IMPACTO DIRETO", 10, COLOR_GRAY_TEXT, True, 14),
        ("Perda de timing = perda de venda", 12, COLOR_RED, True, 6),
    ])

    # Card direito — oportunidade
    add_card(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.7))
    add_text_block(slide, Inches(7.0), Inches(1.65), Inches(5.3), Inches(4.4), [
        ("OPORTUNIDADE IDENTIFICADA", 10, COLOR_GRAY_TEXT, True, 0),
        ("FPS — Ferramentas de Penetração de Solo", 11, COLOR_YELLOW_JD, True, 10),
        ("Desgaste a cada 200h · 4 a 6 compras/ano por máquina", 9.5, COLOR_WHITE, False, 4),
        ("Material Rodante — Tratores de Esteira", 11, COLOR_YELLOW_JD, True, 12),
        ("Desgaste a cada 1.500h · ~1 vez a cada 18 meses", 9.5, COLOR_WHITE, False, 4),
        ("Material Rodante — Escavadeiras", 11, COLOR_YELLOW_JD, True, 12),
        ("Desgaste a cada 3.000h · ~1 vez a cada 3 anos", 9.5, COLOR_WHITE, False, 4),
        ("Padrão previsível = oportunidade estruturada", 10.5, COLOR_GREEN, True, 14),
    ])

    add_callout(slide, "[INSIGHT] Cada hora sem contato é uma janela aberta para a concorrência. O desgaste é previsível — o processo comercial também pode ser.")


def build_slide_solucao(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "motor de leads // visão geral")
    create_slide_title(slide, "MOTOR AUTOMATIZADO DE GERAÇÃO DE LEADS")

    # 5 cards de processo em linha
    steps = [
        ("01", "MONITORAMENTO\nDIÁRIO", "Horímetros de toda a frota lidos e processados todo dia"),
        ("02", "ALERTA\nSEMANAL", "Toda segunda-feira os leads do ciclo são publicados na planilha OneDrive"),
        ("03", "CONTATO\nATIVO", "Consultor/CSA registra status: Venda, Venda Perdida ou Sem Contato"),
        ("04", "DAILY\nREPORT", "Gerentes recebem diariamente os 5 KPIs da campanha por e-mail"),
        ("05", "AUDITORIA\nPROTHEUS", "Toda 'Venda' declarada é cruzada com proposta real no ERP"),
    ]

    card_w = Inches(2.2)
    card_h = Inches(4.5)
    for i, (num, title, desc) in enumerate(steps):
        left = Inches(0.7 + i * 2.42)
        card = add_card(slide, left, Inches(1.55), card_w, card_h)

        # Número grande
        tb_num = slide.shapes.add_textbox(left + Inches(0.1), Inches(1.7), card_w - Inches(0.2), Inches(0.9))
        p = tb_num.text_frame.paragraphs[0]
        p.text = num
        p.font.name = "Segoe UI"
        p.font.size = Pt(36)
        p.font.color.rgb = COLOR_YELLOW_JD
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Título
        tb_t = slide.shapes.add_textbox(left + Inches(0.1), Inches(2.65), card_w - Inches(0.2), Inches(0.8))
        p2 = tb_t.text_frame.paragraphs[0]
        p2.text = title
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_WHITE
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

        # Descrição
        tb_d = slide.shapes.add_textbox(left + Inches(0.12), Inches(3.45), card_w - Inches(0.24), Inches(2.4))
        tb_d.text_frame.word_wrap = True
        p3 = tb_d.text_frame.paragraphs[0]
        p3.text = desc
        p3.font.name = "Segoe UI"
        p3.font.size = Pt(9)
        p3.font.color.rgb = COLOR_GRAY_TEXT
        p3.alignment = PP_ALIGN.CENTER

    add_callout(slide, "[RESULTADO] Processo reativo → processo previsível e auditável. Nenhuma máquina da base ativa some do radar comercial.")


def build_slide_alertas(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "régua de alertas // horímetros")
    create_slide_title(slide, "QUANDO O SISTEMA DISPARA O ALERTA")

    # Tabela de réguas
    add_card(slide, Inches(0.8), Inches(1.5), Inches(7.2), Inches(3.0))
    table = slide.shapes.add_table(4, 3, Inches(0.9), Inches(1.65), Inches(7.0), Inches(2.75)).table

    headers = ["Tipo de Peça", "Equipamento", "Gatilho"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(10, 10, 10)
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Segoe UI"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_YELLOW_JD

    rows = [
        ("FPS (dentes, lâminas, pontas)", "Toda a frota ativa", "A cada +200h de operação"),
        ("Material Rodante", "Tratores de Esteira\n(700J, 750J, 850J, 1050K)", "A cada +1.500h de operação"),
        ("Material Rodante", "Escavadeiras\n(130G, 160G, 200G, 210G, 350G…)", "A cada +3.000h de operação"),
    ]
    for r, (tp, eq, gt) in enumerate(rows):
        colors = [COLOR_YELLOW_JD, COLOR_WHITE, COLOR_GREEN]
        for c, (txt, col) in enumerate(zip([tp, eq, gt], colors)):
            cell = table.cell(r + 1, c)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_GRAY_DARK
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Segoe UI"
            p.font.size = Pt(9.5)
            p.font.color.rgb = col

    # Card ciclo de vida (direita)
    add_card(slide, Inches(8.2), Inches(1.5), Inches(4.7), Inches(4.6))
    add_text_block(slide, Inches(8.4), Inches(1.65), Inches(4.3), Inches(4.3), [
        ("CICLO DE VIDA DO LEAD", 10, COLOR_GRAY_TEXT, True, 0),
        ("Máquina atinge o limiar de horas", 9.5, COLOR_WHITE, False, 10),
        ("↓  Sistema gera o alerta na planilha", 9.5, COLOR_YELLOW_JD, True, 4),
        ("↓  Consultor faz o contato", 9.5, COLOR_WHITE, False, 4),
        ("VENDA ou VENDA PERDIDA", 10, COLOR_GREEN, True, 10),
        ("→ Marco zero atualizado\n→ Novo ciclo inicia", 9, COLOR_WHITE, False, 2),
        ("SEM CONTATO / EM BRANCO", 10, COLOR_AMBER, True, 10),
        ("→ Lead permanece na fila\n→ Aging cresce no Daily Report", 9, COLOR_WHITE, False, 2),
    ])

    # Bloco info adicional
    add_card(slide, Inches(0.8), Inches(4.7), Inches(7.2), Inches(1.45))
    add_text_block(slide, Inches(1.0), Inches(4.85), Inches(6.8), Inches(1.25), [
        ("POR QUE LIMIARES DIFERENTES PARA RODANTE?", 9.5, COLOR_GRAY_TEXT, True, 0),
        ("Tratores de esteira se locomovem continuamente — esteiras em atrito constante com o solo. "
         "Escavadeiras trabalham paradas e rodam as esteiras apenas em transferências. "
         "Por isso o desgaste do trator é ~2× mais rápido.", 9, COLOR_WHITE, False, 4),
    ])

    add_callout(slide, "[REGRA] Após Venda ou Venda Perdida, o horímetro é zerado e o ciclo reinicia. Leads Sem Contato continuam acumulando horas.")


def build_slide_carga_inicial(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "régua de alertas // carga inicial")
    create_slide_title(slide, "PRIMEIRA CARGA: PARTINDO DO ZERO SEM RUÍDO")

    # Problema da carga inicial
    add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(4.7))
    add_text_block(slide, Inches(1.0), Inches(1.65), Inches(5.4), Inches(4.4), [
        ("O DESAFIO DA CARGA INICIAL", 10, COLOR_GRAY_TEXT, True, 0),
        ("Na primeira vez que o sistema processa um chassis, não há histórico "
         "de quando a última peça foi trocada.", 10, COLOR_WHITE, False, 8),
        ("Regra aplicada:", 10, COLOR_YELLOW_JD, True, 12),
        ("O horímetro atual é registrado como MARCO ZERO.", 13, COLOR_WHITE, True, 6),
        ("O primeiro alerta só é gerado após acumular as horas da régua a partir desse momento.", 10, COLOR_WHITE, False, 6),
        ("Por quê?", 10, COLOR_YELLOW_JD, True, 12),
        ("Evita gerar centenas de alertas simultâneos ao ativar a campanha — "
         "o que sobrecarregaria a equipe no dia 1.", 9.5, COLOR_GRAY_TEXT, False, 4),
    ])

    # Exemplo prático
    add_card(slide, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.7))
    add_text_block(slide, Inches(7.1), Inches(1.65), Inches(5.2), Inches(4.4), [
        ("EXEMPLO PRÁTICO — ESCAVADEIRA 200G", 10, COLOR_GRAY_TEXT, True, 0),
        ("Horímetro no momento da ativação:", 9.5, COLOR_WHITE, False, 10),
        ("4.200 horas", 22, COLOR_YELLOW_JD, True, 4),
        ("→ Marco zero registrado: 4.200h", 10, COLOR_WHITE, False, 8),
        ("Próximo alerta de FPS:", 9.5, COLOR_GRAY_TEXT, False, 10),
        ("4.400h  (+200h)", 14, COLOR_GREEN, True, 2),
        ("Próximo alerta de Rodante:", 9.5, COLOR_GRAY_TEXT, False, 8),
        ("7.200h  (+3.000h)", 14, COLOR_AMBER, True, 2),
        ("A campanha começa de forma ordenada,\nsem ruído desnecessário.", 9.5, COLOR_GRAY_TEXT, False, 10),
    ])

    add_callout(slide, "[GARANTIA] A ativação da campanha é suave. Os alertas chegam de forma distribuída ao longo das semanas — nunca todos de uma vez.")


def build_slide_kpi_bloco1(prs: Presentation, has_screenshot: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "indicadores // desempenho")
    create_slide_title(slide, "INDICADORES DE DESEMPENHO (KPIs)")

    if has_screenshot and KPI_SCREENSHOT.exists():
        # Screenshot recortado para mostrar só os cards KPI (parte superior)
        slide.shapes.add_picture(str(KPI_SCREENSHOT), Inches(0.8), Inches(1.5), Inches(6.2), Inches(3.5))

    # Anotações com setas textuais (direita)
    add_card(slide, Inches(7.2), Inches(1.5), Inches(5.4), Inches(4.7))
    add_text_block(slide, Inches(7.4), Inches(1.65), Inches(5.0), Inches(4.4), [
        ("OS 4 INDICADORES PRIMÁRIOS", 10, COLOR_GRAY_TEXT, True, 0),
        ("◀  ADESÃO COMERCIAL", 10.5, COLOR_YELLOW_JD, True, 10),
        ("% dos leads com feedback registrado pela equipe.\nMede velocidade e engajamento no contato.", 9, COLOR_WHITE, False, 2),
        ("◀  CONVERSÃO ACUMULADA", 10.5, COLOR_YELLOW_JD, True, 10),
        ("% de leads tratados que resultaram em venda.\nMede eficiência do discurso comercial.", 9, COLOR_WHITE, False, 2),
        ("◀  ADERÊNCIA DE PROPOSTAS", 10.5, COLOR_YELLOW_JD, True, 10),
        ("% de 'Vendas' com proposta real no Protheus.\nA Ponte da Verdade — elimina o autodeclaratório.", 9, COLOR_WHITE, False, 2),
        ("◀  FATURAMENTO REALIZADO", 10.5, COLOR_GREEN, True, 10),
        ("Valor em R$ das propostas validadas no ERP.", 9, COLOR_WHITE, False, 2),
    ])

    add_callout(slide, "[GOVERNANÇA] Os KPIs são enviados diariamente por e-mail para os 8 destinatários do Daily Report — sem necessidade de acessar planilha.")


def build_slide_kpi_bloco2(prs: Presentation, has_screenshot: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "indicadores // metodologia")
    create_slide_title(slide, "METODOLOGIA E AGING COMERCIAL")

    if has_screenshot and KPI_SCREENSHOT.exists():
        # Screenshot recortado para mostrar a tabela de aging (parte inferior do HTML)
        slide.shapes.add_picture(str(KPI_SCREENSHOT), Inches(0.8), Inches(1.5), Inches(6.2), Inches(3.5))

    # Anotações (direita)
    add_card(slide, Inches(7.2), Inches(1.5), Inches(5.4), Inches(4.7))
    add_text_block(slide, Inches(7.4), Inches(1.65), Inches(5.0), Inches(4.4), [
        ("AGING E PIPELINE FINANCEIRO", 10, COLOR_GRAY_TEXT, True, 0),
        ("◀  AGING DO LEAD", 10.5, COLOR_YELLOW_JD, True, 10),
        ("Dias médios sem primeiro contato após o alerta.\nAlerta para perda de timing preventivo.", 9, COLOR_WHITE, False, 2),
        ("◀  PIPELINE FINANCEIRO", 10.5, COLOR_GREEN, True, 10),
        ("Soma em R$ das propostas abertas no Protheus\nvinculadas à campanha. Visibilidade financeira\ndo que está em negociação.", 9, COLOR_WHITE, False, 2),
        ("LEITURA DO PAINEL", 10, COLOR_GRAY_TEXT, True, 14),
        ("Aging alto = timing preventivo sendo perdido", 9.5, COLOR_RED, True, 4),
        ("Aderência baixa = resultado autodeclaratório", 9.5, COLOR_AMBER, True, 4),
        ("Conversão crescente = discurso eficaz", 9.5, COLOR_GREEN, True, 4),
    ])

    add_callout(slide, "[AUTONOMIA] Gerentes e coordenadores recebem os mesmos KPIs diariamente — mesma visibilidade, comparação justa entre regiões.")


def build_slide_governanca(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "governança // auditoria")
    create_slide_title(slide, "A PONTE DA VERDADE: PLANILHA × PROTHEUS")

    # Fluxo esquerda
    add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(4.7))
    add_text_block(slide, Inches(1.0), Inches(1.65), Inches(5.4), Inches(4.4), [
        ("COMO A AUDITORIA FUNCIONA", 10, COLOR_GRAY_TEXT, True, 0),
        ("1. Consultor registra 'Venda' na planilha", 10, COLOR_WHITE, False, 10),
        ("2. Sistema extrai leads marcados como Venda", 10, COLOR_WHITE, False, 6),
        ("3. Cruza com propostas reais do ERP Protheus", 10, COLOR_WHITE, False, 6),
        ("4. Calcula Aderência de Propostas (KPI)", 10, COLOR_WHITE, False, 6),
        ("", 6, COLOR_WHITE, False, 2),
        ("Venda SEM proposta no Protheus:", 10, COLOR_AMBER, True, 8),
        ("→ Aparece imediatamente no KPI de Aderência", 9.5, COLOR_WHITE, False, 2),
        ("→ Fica destacado no Daily Report para o gerente", 9.5, COLOR_WHITE, False, 2),
        ("Venda COM proposta no Protheus:", 10, COLOR_GREEN, True, 8),
        ("→ Computada no Pipeline Financeiro em R$", 9.5, COLOR_WHITE, False, 2),
    ])

    # Proteções de dados (direita)
    add_card(slide, Inches(6.9), Inches(1.5), Inches(5.6), Inches(4.7))
    add_text_block(slide, Inches(7.1), Inches(1.65), Inches(5.2), Inches(4.4), [
        ("PROTEÇÃO DOS DADOS NA PLANILHA", 10, COLOR_GRAY_TEXT, True, 0),
        ("Colunas de origem BLOQUEADAS:", 10, COLOR_YELLOW_JD, True, 10),
        ("• Chassi, Cliente, CNPJ, Modelo\n• Motivo do alerta, Horímetro", 9.5, COLOR_WHITE, False, 4),
        ("Consultores só editam:", 10, COLOR_YELLOW_JD, True, 10),
        ("• Status (Venda / Venda Perdida / Sem Contato)\n• Campo de observações livres", 9.5, COLOR_WHITE, False, 4),
        ("Snapshot diário:", 10, COLOR_YELLOW_JD, True, 10),
        ("Antes de cada atualização semanal, o sistema\nsalva o estado completo da planilha para\nrastreabilidade histórica.", 9, COLOR_WHITE, False, 4),
    ])

    add_callout(slide, "[FIM DO AUTODECLARATÓRIO] Qualquer resultado reportado sem evidência no ERP é detectado automaticamente — tornando os números confiáveis.")


def build_slide_responsabilidades(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "papéis e responsabilidades")
    create_slide_title(slide, "QUEM FAZ O QUÊ NA CAMPANHA")

    add_card(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.7))

    table = slide.shapes.add_table(6, 3, Inches(0.9), Inches(1.65), Inches(11.5), Inches(4.5)).table

    headers = ["Papel", "Responsável", "Ação na Campanha"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(10, 10, 10)
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Segoe UI"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_YELLOW_JD

    rows = [
        ("Gerente de Peças", "Roberto Reis",
         "Lidera reunião semanal de revisão de metas (terças); avalia KPIs consolidados"),
        ("Gerentes / Coordenadores Regionais", "Pedro Sarnaglia · Leandro Silva\nMarcelo Costa · Luciana Borges",
         "Acompanham o Daily Report diário; cobram aging elevado dos consultores; gerenciam equipe regional"),
        ("CSA", "Murilo Nunes",
         "Realiza contato ativo com clientes da carteira; registra status na planilha"),
        ("Consultores de Vendas", "A definir (mapeamento Murilo)",
         "Realizam o contato ativo com os clientes; registram Venda / Venda Perdida / Sem Contato"),
        ("Engenharia de Dados", "Victor Bernardi",
         "Mantém o motor de cálculo, atualiza a planilha semanalmente e dispara o Daily Report"),
    ]
    for r, (papel, resp, acao) in enumerate(rows):
        row_data = [papel, resp, acao]
        for c, txt in enumerate(row_data):
            cell = table.cell(r + 1, c)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_GRAY_DARK
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Segoe UI"
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_YELLOW_JD if c == 0 else COLOR_WHITE

    add_callout(slide, "[PENDÊNCIA] Segmentação de clientes por CSA/consultor em definição. Assim que Murilo enviar o mapeamento, a planilha será configurada com Responsável automático.")


def build_slide_proximos_passos(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_theme(prs, slide)
    create_eyebrow(slide, "roadmap // implantação")
    create_slide_title(slide, "PRIMEIROS MARCOS DA CAMPANHA")

    marcos = [
        ("SEMANA 1", "Primeiro ciclo semanal de alertas publicado na planilha OneDrive", COLOR_YELLOW_JD),
        ("SEMANA 1", "Primeiro Daily Report enviado para todos os 8 destinatários", COLOR_YELLOW_JD),
        ("MÊS 1", "Primeira auditoria de Aderência de Propostas: planilha × Protheus", COLOR_AMBER),
        ("MÊS 2", "Primeiro benchmark de conversão por região disponível", COLOR_GREEN),
    ]

    card_w = Inches(2.7)
    for i, (periodo, desc, color) in enumerate(marcos):
        left = Inches(0.7 + i * 3.0)
        add_card(slide, left, Inches(1.55), card_w, Inches(4.5))

        # Período
        tb_p = slide.shapes.add_textbox(left + Inches(0.1), Inches(1.75), card_w - Inches(0.2), Inches(0.6))
        p = tb_p.text_frame.paragraphs[0]
        p.text = periodo
        p.font.name = "Segoe UI"
        p.font.size = Pt(20)
        p.font.color.rgb = color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Linha decorativa
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.5), Inches(2.4), card_w - Inches(1.0), Inches(0.04))
        ln.fill.solid()
        ln.fill.fore_color.rgb = color
        ln.line.fill.background()

        # Descrição
        tb_d = slide.shapes.add_textbox(left + Inches(0.15), Inches(2.55), card_w - Inches(0.3), Inches(3.3))
        tb_d.text_frame.word_wrap = True
        p2 = tb_d.text_frame.paragraphs[0]
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_WHITE
        p2.alignment = PP_ALIGN.CENTER

    add_callout(slide, "[OBJETIVO] Da semana 1 ao mês 2: do primeiro alerta ao primeiro benchmark de resultado — campanha rodando e auditável em 60 dias.")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    print("=== GERANDO APRESENTAÇÃO LEADS PREVENTIVOS ===")

    has_screenshot = capture_kpi_screenshot()

    print(" -> Instanciando apresentação (13.33\" × 7.5\")...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print(" -> Slide 1: Capa")
    build_slide_capa(prs)
    print(" -> Slide 2: O Problema")
    build_slide_problema(prs)
    print(" -> Slide 3: A Solução")
    build_slide_solucao(prs)
    print(" -> Slide 4: Como os Alertas São Gerados")
    build_slide_alertas(prs)
    print(" -> Slide 5: Carga Inicial vs Ciclo Contínuo")
    build_slide_carga_inicial(prs)
    print(" -> Slide 6: KPI Bloco 1")
    build_slide_kpi_bloco1(prs, has_screenshot)
    print(" -> Slide 7: KPI Bloco 2")
    build_slide_kpi_bloco2(prs, has_screenshot)
    print(" -> Slide 8: Governança — Ponte da Verdade")
    build_slide_governanca(prs)
    print(" -> Slide 9: Responsabilidades")
    build_slide_responsabilidades(prs)
    print(" -> Slide 10: Próximos Passos")
    build_slide_proximos_passos(prs)

    prs.save(str(OUTPUT_PPTX))
    size = OUTPUT_PPTX.stat().st_size if OUTPUT_PPTX.exists() else 0
    print(f"\n[SUCESSO] {OUTPUT_PPTX.name} ({size:,} bytes)")
    print("==============================================")


if __name__ == "__main__":
    main()
